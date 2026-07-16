"""Rebuild v2 deck: site-license pricing.
- Slide 8: reword per-surgeon -> per-site subscription
- Insert Financial Projections (11) + Long View (12) slides
- Renumber old Ask slide 11 -> 13"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

TEAL = RGBColor(0x14, 0xB8, 0xA6)
NAVY = RGBColor(0x0E, 0x22, 0x3B)
BODY = RGBColor(0x24, 0x2E, 0x3A)
GRAY = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
FONT = "DejaVu Sans"

src = "/Users/nathanaelguitar/Downloads/AlphaAgent_Pitch_Deck.pptx"
out = "/Users/nathanaelguitar/Downloads/AlphaAgent_Pitch_Deck_v2.pptx"
p = Presentation(src)

# ---- Slide 8: per-surgeon -> per-site wording ----
swaps = {
    "▸  Recurring per-surgeon subscription to hospitals, surgical groups, and ASCs.":
    "▸  Site license: $5–20K/mo per hospital or ASC, paid by the institution.",
    "↓   per-surgeon subscription": "↓   site subscription",
}
for shape in p.slides[7].shapes:
    if not shape.has_text_frame: continue
    for para in shape.text_frame.paragraphs:
        full = "".join(r.text for r in para.runs)
        if full in swaps:
            para.runs[0].text = swaps[full]
            for r in para.runs[1:]: r.text = ""

def solid(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def txt(slide, l, t, w, h, lines, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (s, sz, b, c, sa) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align; para.space_after = Pt(sa)
        run = para.add_run(); run.text = s
        f = run.font; f.name = FONT; f.size = Pt(sz); f.bold = b; f.color.rgb = c
    return tb

def chrome(slide, kicker, title, page):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, p.slide_width, Inches(0.09))
    solid(bar, TEAL)
    txt(slide, 0.7, 0.42, 11.93, 0.32, [(kicker, 12, True, TEAL, 0)])
    txt(slide, 0.7, 0.78, 11.93, 0.6, [(title, 25, True, NAVY, 0)])
    ul = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.78), Inches(1.5), Inches(0.042))
    solid(ul, TEAL)
    txt(slide, 0.7, 7.08, 6.0, 0.3, [("AlphaAgent Healthcare Solutions  ·  Confidential", 9, False, GRAY, 0)])
    txt(slide, 12.13, 7.08, 0.5, 0.3, [(str(page), 9, False, GRAY, 0)])

def table(slide, data, l, t, w, col0_w, row_h=0.36, hi_row=None, hi_cell=None):
    """data[0] is header. hi_row: teal-value row. hi_cell: (i,j) extra teal bold."""
    n_rows, n_cols = len(data), len(data[0])
    tbl = slide.shapes.add_table(n_rows, n_cols, Inches(l), Inches(t),
                                 Inches(w), Inches(row_h * n_rows)).table
    tbl.columns[0].width = Inches(col0_w)
    cw = (w - col0_w) / (n_cols - 1)
    for j in range(1, n_cols): tbl.columns[j].width = Inches(cw)
    for i, row in enumerate(data):
        tbl.rows[i].height = Inches(row_h)
        for j, v in enumerate(row):
            cell = tbl.cell(i, j); cell.text = v
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if i == 0 else (LIGHT if i % 2 == 0 else WHITE)
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.RIGHT if j else PP_ALIGN.LEFT
            teal = (i == hi_row and j > 0) or ((i, j) == hi_cell)
            for r in (para.runs or [para.add_run()]):
                r.font.name = FONT; r.font.size = Pt(12)
                r.font.bold = (i == 0) or (j == 0) or teal or i == hi_row
                r.font.color.rgb = WHITE if i == 0 else (TEAL if teal else BODY)
    return tbl

def panel(slide, header, items):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(8.55), Inches(2.05), Inches(4.08), Inches(4.75))
    card.adjustments[0] = 0.045
    solid(card, NAVY)
    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.55), Inches(2.05), Inches(4.08), Inches(0.06))
    solid(strip, TEAL)
    txt(slide, 8.87, 2.32, 3.44, 0.32, [(header, 11, True, TEAL, 0)])
    txt(slide, 8.87, 2.78, 3.44, 3.9, items)

# ================= Slide 11: Financial projections =================
s = p.slides.add_slide(p.slide_layouts[6])
chrome(s, "FINANCIAL PROJECTIONS", "Cashflow-positive in Year 3 — on $2M total capital", 11)
table(s, [
    ("", "Year 1", "Year 2", "Year 3"),
    ("Hospital sites (EOY)", "3", "10", "25"),
    ("Exit ARR", "$360K", "$1.2M", "$3.0M"),
    ("Revenue", "$90K", "$860K", "$2.21M"),
    ("Gross margin", "23%", "74%", "76%"),
    ("Opex", "$348K", "$870K", "$1.50M"),
    ("Net income", "($327K)", "($236K)", "+$189K"),
], l=0.7, t=2.15, w=7.4, col0_w=3.0, hi_row=2, hi_cell=(6, 3))
txt(s, 0.7, 5.05, 7.5, 1.7, [
    ("▸  Bottom-up: $10K/mo blended site license ($5K entry service line → $15–20K expanded).", 11.5, False, BODY, 6),
    ("▸  8-month pilots (procurement + IT review) convert to paid in month 9; Y3 ramp follows 510(k).", 11.5, False, BODY, 6),
    ("▸  $2M total capital to breakeven: this $500K pre-seed + ~$1.5M seed raised on converted pilots.", 11.5, False, BODY, 0),
])
panel(s, "UNIT ECONOMICS / SITE", [
    ("$10K / month site license (blended)", 12, True, WHITE, 2),
    ("$1K / month cloud, inference & support", 12, False, WHITE, 12),
    ("$20K hardware — 10 glasses, free to the hospital", 12, True, WHITE, 2),
    ("recovered in 2.2 months of gross profit", 12, False, WHITE, 12),
    ("90% software gross margin at steady state", 12, False, WHITE, 12),
    ("Every site = $120K ARR against a one-time $20K hardware cost.", 12, True, TEAL, 0),
])

# ================= Slide 12: The long view =================
s2 = p.slides.add_slide(p.slide_layouts[6])
chrome(s2, "THE LONG VIEW", "$135M ARR at 4% of US surgical facilities", 12)
table(s2, [
    ("", "Year 3", "Year 5", "Year 7", "Year 10"),
    ("Sites live", "25", "90", "260", "750"),
    ("Avg ACV / site", "$120K", "$150K", "$165K", "$180K"),
    ("ARR", "$3.0M", "$13.5M", "$42.9M", "$135M"),
    ("US penetration", "0.1%", "0.5%", "1.5%", "4.4%"),
], l=0.7, t=2.15, w=7.4, col0_w=2.6, hi_row=3)
txt(s2, 0.7, 4.35, 7.5, 2.4, [
    ("▸  Market: ~6,100 US hospitals + ~11,000 ASCs. 750 sites in ten years is 4.4% penetration.", 11.5, False, BODY, 6),
    ("▸  ACV compounds inside each site: one service line → department → system standard of care.", 11.5, False, BODY, 6),
    ("▸  Not shown: the second act — every approved record trains the proprietary multimodal surgical dataset, and diagnostics licensing stacks on the subscription base.", 11.5, False, BODY, 0),
])
panel(s2, "WHY THIS COMPOUNDS", [
    ("Sticky — ripping it out means going back to the keyboard and the phone.", 12, False, WHITE, 12),
    ("Regulated — every 510(k) quarter we bank becomes a barrier behind us.", 12, False, WHITE, 12),
    ("Self-funding — breakeven in Year 3 means growth capital scales the curve, it does not keep the lights on.", 12, False, WHITE, 12),
    ("Data moat — 750 sites of labeled surgical video is a dataset nobody can buy.", 12, True, TEAL, 0),
])

# move new slides into position 11 and 12
sldIdLst = p.slides._sldIdLst
ids = list(sldIdLst)
sldIdLst.remove(ids[-2]); sldIdLst.insert(10, ids[-2])
sldIdLst.remove(ids[-1]); sldIdLst.insert(11, ids[-1])

# renumber old Ask slide footer 11 -> 13
for shape in p.slides[12].shapes:
    if shape.has_text_frame and shape.text_frame.text.strip() == "11":
        shape.text_frame.paragraphs[0].runs[0].text = "13"

p.save(out)
print("saved", out, "slides:", len(p.slides._sldIdLst))
