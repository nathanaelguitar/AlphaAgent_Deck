#!/usr/bin/env python3
"""
Build AlphaAgent investor deck as a .pptx (16:9). Content mirrors pitch_deck.md.
Designed for DejaVu Sans metrics (the font available on this machine), with a
strict layout grid so nothing overlaps. Speaker notes attached to every slide.

Run:  ./.venv/bin/python build_pptx.py   ->   AlphaAgent_Pitch_Deck.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- Palette ---------------------------------------------------------------
NAVY   = RGBColor(0x0E, 0x22, 0x3B)
NAVY2  = RGBColor(0x13, 0x2C, 0x4A)   # panel on dark
TEAL   = RGBColor(0x14, 0xB8, 0xA6)
LIGHT  = RGBColor(0xF1, 0xF5, 0xF9)
LINE   = RGBColor(0xDD, 0xE4, 0xEB)
INK    = RGBColor(0x24, 0x2E, 0x3A)
MUTE   = RGBColor(0x64, 0x74, 0x8B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GOLD   = RGBColor(0xF5, 0xC0, 0x5A)
SOFT   = RGBColor(0x9F, 0xB3, 0xC8)   # muted text on navy

FONT = "DejaVu Sans"

SW, SH = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.7)

prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH
BLANK = prs.slide_layouts[6]


# ---- Primitives --------------------------------------------------------------
def _style(run, size, color, bold=False, italic=False):
    f = run.font
    f.size, f.bold, f.italic, f.name = Pt(size), bold, italic, FONT
    f.color.rgb = color


def shape(slide, kind, x, y, w, h, fill=None, line_color=None, line_w=0.75, round_=None):
    sp = slide.shapes.add_shape(kind, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line_color is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line_color; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if round_ is not None and sp.adjustments:
        try: sp.adjustments[0] = round_
        except Exception: pass
    return sp


def rect(slide, x, y, w, h, fill, line_color=None):
    return shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill, line_color)


def card(slide, x, y, w, h, fill, line_color=None):
    return shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, line_color, round_=0.045)


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is a list of (txt,size,color,bold,italic)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        if line_spacing: p.line_spacing = line_spacing
        for (t, sz, col, *rest) in para:
            bold = rest[0] if len(rest) > 0 else False
            ital = rest[1] if len(rest) > 1 else False
            r = p.add_run(); r.text = t; _style(r, sz, col, bold, ital)
    return tb


def notes(slide, s):
    slide.notes_slide.notes_text_frame.text = s


def footer(slide, n, dark=False):
    col = SOFT if dark else MUTE
    text(slide, MARGIN, Inches(7.08), Inches(6), Inches(0.3),
         [[("AlphaAgent Healthcare Solutions  ·  Confidential", 9, col)]])
    text(slide, SW - Inches(1.2), Inches(7.08), Inches(0.5), Inches(0.3),
         [[(str(n), 9, col)]], align=PP_ALIGN.RIGHT)


def header(slide, kicker, title):
    """Header zone: y 0.45–1.85. Content must start at y >= 2.05."""
    rect(slide, 0, 0, SW, Inches(0.09), TEAL)
    text(slide, MARGIN, Inches(0.42), SW - 2*MARGIN, Inches(0.32),
         [[(kicker.upper(), 12, TEAL, True)]])
    text(slide, MARGIN, Inches(0.78), SW - 2*MARGIN, Inches(0.95),
         [[(title, 25, NAVY, True)]], line_spacing=1.05)
    rect(slide, MARGIN, Inches(1.78), Inches(1.5), Pt(3), GOLD)


def bullet_list(slide, items, x, y, w, size=14, gap=12, dark=False):
    """items: (text, bold) or (text, bold, sub=True)."""
    tb = slide.shapes.add_textbox(x, y, w, Inches(0.4))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    body = WHITE if dark else INK
    sub_c = SOFT if dark else MUTE
    for i, it in enumerate(items):
        t = it[0]; bold = it[1] if len(it) > 1 else False
        sub = it[2] if len(it) > 2 else False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.12
        r = p.add_run(); r.text = ("–  " if sub else "▸  ") + t
        _style(r, size - (1 if sub else 0), sub_c if sub else body, bold)
        if not sub:
            # teal marker for top-level bullets
            r.font.color.rgb = body
    return tb


def panel(slide, title, lines, x=Inches(8.55), y=Inches(2.05), w=Inches(4.08), h=Inches(4.75)):
    """Right-hand callout card on light slides."""
    card(slide, x, y, w, h, NAVY)
    rect(slide, x, y, w, Inches(0.06), TEAL)
    inner_x, inner_w = x + Inches(0.32), w - Inches(0.64)
    text(slide, inner_x, y + Inches(0.3), inner_w, Inches(0.35),
         [[(title.upper(), 11, TEAL, True)]])
    tb = slide.shapes.add_textbox(inner_x, y + Inches(0.78), inner_w, h - Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, ln in enumerate(lines):
        t, hl = (ln, False) if isinstance(ln, str) else ln
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(9); p.line_spacing = 1.15
        r = p.add_run(); r.text = t
        _style(r, 12.5, GOLD if hl else WHITE, hl)
    return tb


# ===========================================================================
# SLIDE 1 — Title
# ===========================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, 0, SW, Inches(0.09), TEAL)
text(s, MARGIN, Inches(1.55), Inches(11.9), Inches(0.4),
     [[("MEDTECH  ·  AGENTIC AI  ·  SEED", 13, TEAL, True)]])
text(s, MARGIN, Inches(2.05), Inches(11.9), Inches(1.9),
     [[("AlphaAgent", 48, WHITE, True)],
      [("Healthcare Solutions", 48, WHITE, True)]], line_spacing=1.0, space_after=2)
rect(s, MARGIN, Inches(4.15), Inches(2.2), Pt(3), GOLD)
text(s, MARGIN, Inches(4.42), Inches(11.9), Inches(0.55),
     [[("The autonomous documentation layer for the operating room.", 19, TEAL)]])
text(s, MARGIN, Inches(5.12), Inches(11.9), Inches(0.5),
     [[("Smart glasses + an agentic AI that documents surgery, hands-free.", 14, SOFT)]])
text(s, MARGIN, Inches(6.55), Inches(11.9), Inches(0.4),
     [[("Nathanael Gill, CEO    ·    Manny Figueroa, CTO    ·    Confidential — July 2026", 11, SOFT)]])
notes(s, "AlphaAgent puts an autonomous AI agent behind a surgeon's eyes. It watches, listens, and "
"writes the clinical record so the surgeon can stop being a data-entry clerk. We pair enterprise "
"smart glasses with an agent that doesn't just transcribe, it acts: it populates the EHR and waits "
"for a one-word voice approval.")

# ===========================================================================
# SLIDE 2 — Problem
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "The Problem", "Surgeons are documentation clerks with scalpels")
bullet_list(s, [
    ("Documentation is the #1 driver of clinician burnout — hours lost every day.", True),
    ("Dual-task interference (care + charting) drives up error rates.",),
    ("Pre/post-op photos are taken on personal phones —", False),
    ("a contamination and compliance risk inside a sterile field.", False, True),
    ("Today's AI scribes stop at a text note — a human still does the clicking.",),
], MARGIN, Inches(2.15), Inches(7.4), size=15, gap=14)
panel(s, "The two failures", [
    ("The keyboard:", True),
    "every minute typing is a minute away from the patient.",
    ("The phone:", True),
    "a personal device entering the sterile field for wound photos.",
    "Both are in the wrong place. We remove both.",
])
footer(s, 2)
notes(s, "Two things are broken. First, the surgeon is drowning in charting; every minute on the "
"keyboard is a minute not on the patient. Second, the visual record gets captured on a personal "
"phone that has no business in a sterile field. Ambient scribes help with words but never touch the "
"images, and they don't act. We remove the keyboard and the phone entirely.")

# ===========================================================================
# SLIDE 3 — Solution (3 step cards)
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "The Solution", "Capture → Reason → Act. Hands-free.")
steps = [
    ("1 · CAPTURE", "Vuzix LX1 glasses capture 12MP photos + OR audio — hands-free, sterile."),
    ("2 · REASON", "A cloud agent with persistent memory turns capture into a structured clinical record."),
    ("3 · ACT", "The agent writes EHR fields via FHIR. The surgeon approves or edits by voice."),
]
cw, ch, gap_w = Inches(3.85), Inches(3.0), Inches(0.27)
for i, (t, d) in enumerate(steps):
    x = MARGIN + i * (cw + gap_w)
    card(s, x, Inches(2.15), cw, ch, LIGHT, LINE)
    rect(s, x, Inches(2.15), cw, Inches(0.07), TEAL)
    text(s, x + Inches(0.3), Inches(2.5), cw - Inches(0.6), Inches(0.4),
         [[(t, 14, NAVY, True)]])
    text(s, x + Inches(0.3), Inches(3.05), cw - Inches(0.6), Inches(1.9),
         [[(d, 13, INK)]], line_spacing=1.25)
text(s, MARGIN, Inches(5.55), Inches(12.0), Inches(0.6),
     [[("Zero phones. Zero keyboards. ", 16, NAVY, True),
       ("The surgeon stays with the patient — “approve” finishes the chart.", 16, TEAL, True)]])
footer(s, 3)
notes(s, "Three steps, none requiring the surgeon's hands. The glasses capture the moment. Our agent "
"— with memory and tool-calling — turns it into a real clinical record and writes it into the EHR "
"via FHIR. The surgeon just says 'approve' or dictates an edit. Documentation is done before they've "
"scrubbed out.")

# ===========================================================================
# SLIDE 4 — Product & Moat
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Product & Technical Moat", "An agent that acts, on hardware that computes")
bullet_list(s, [
    ("On-device compute: the LX1 is a full Android 15 computer.", True),
    ("8-core Qualcomm · 6GB/128GB · Wi-Fi 6E — PHI pre-processed before it leaves the room.", False, True),
    ("Agentic core: persistent memory, deterministic tool-calling, FHIR EHR writes.", True),
    ("Proprietary TypeScript/Python stack — not a wrapper around a chat model.", False, True),
    ("Compliance-native: Azure Blob (BAA) for media, GCP Gemini (BAA) for voice.",),
    ("Data flywheel: every approved record becomes labeled training data.", True),
], MARGIN, Inches(2.1), Inches(7.5), size=14, gap=11)
panel(s, "Architecture", [
    "LX1 glasses  (Android 15, on-device)",
    "↓   edge pre-processing",
    "Cloud agent  (memory + tools)",
    "↓   FHIR API",
    "EHR — fields populated, voice-approved",
    ("Every workflow → proprietary surgical dataset.", True),
])
footer(s, 4)
notes(s, "Why this is defensible: the glasses are a computer, so data can be handled on-device first "
"— the first question every hospital CISO asks. Our real IP is the agent: memory, reliable "
"tool-calling, direct EHR integration. Anyone can call an LLM; few can make an agent that reliably "
"acts inside an enterprise workflow. Every surgery generates proprietary labeled data a software- or "
"hardware-only competitor cannot replicate.")

# ===========================================================================
# SLIDE 5 — Why Now
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Why Now", "Three curves crossed at once")
cards5 = [
    ("AGENTS GOT RELIABLE", "AI agents crossed the line from chatting to acting in 2025–26 — reliable enough to trust with a medical record."),
    ("GLASSES GOT VIABLE", "Enterprise eyewear finally has the on-device compute, optics, and mics for the OR."),
    ("HOSPITALS ARE DESPERATE", "Staffing shortages and burnout are at all-time highs — buyers are actively looking."),
]
cw, ch, gap_w = Inches(3.85), Inches(3.1), Inches(0.27)
for i, (t, d) in enumerate(cards5):
    x = MARGIN + i * (cw + gap_w)
    card(s, x, Inches(2.15), cw, ch, NAVY)
    rect(s, x, Inches(2.15), cw, Inches(0.07), TEAL)
    text(s, x + Inches(0.3), Inches(2.5), cw - Inches(0.6), Inches(0.65),
         [[(t, 13, TEAL, True)]], line_spacing=1.1)
    text(s, x + Inches(0.3), Inches(3.2), cw - Inches(0.6), Inches(1.9),
         [[(d, 12.5, WHITE)]], line_spacing=1.25)
text(s, MARGIN, Inches(5.65), Inches(12.0), Inches(0.55),
     [[("This product was technically impossible three years ago. ", 15, NAVY, True),
       ("The window is open — and it won't stay uncontested.", 15, MUTE)]])
footer(s, 5)
notes(s, "Three years ago agents weren't reliable enough to trust with a medical record, and glasses "
"weren't powerful enough to run anything on-device. Both crossed the line in the last 18 months, at "
"the moment hospitals hit a staffing crisis. That convergence is the window, and it won't stay open "
"uncontested.")

# ===========================================================================
# SLIDE 6 — Competitive 2x2
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Competitive Landscape", "We own the quadrant neither side can reach")

mx, my = Inches(1.35), Inches(2.1)
qw, qh, qgap = Inches(2.95), Inches(2.12), Inches(0.14)

def quad(cx, cy, title, sub, fill, tcol, scol, line_c=None):
    c = card(s, cx, cy, qw, qh, fill, line_c)
    tf = c.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.22)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(5)
    r = p.add_run(); r.text = title; _style(r, 13.5, tcol, True)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.line_spacing = 1.12
    r = p2.add_run(); r.text = sub; _style(r, 10.5, scol)

# top-left / top-right
quad(mx, my, "DAX · Abridge · Augmedix",
     "AI scribes — no eyewear, no visual capture, stop at a draft note", WHITE, INK, MUTE, LINE)
quad(mx + qw + qgap, my, "AlphaAgent",
     "Wearable + agentic — autonomously completes the record", TEAL, WHITE, WHITE)
# bottom-left / bottom-right
quad(mx, my + qh + qgap, "Legacy manual charting",
     "Keyboard + personal phone", LIGHT, MUTE, MUTE, LINE)
quad(mx + qw + qgap, my + qh + qgap, "Vuzix (hardware only)",
     "Enterprise eyewear — no clinical agent, not built for the OR", WHITE, INK, MUTE, LINE)

# axes
ax_y = my + 2*qh + qgap + Inches(0.12)
text(s, mx, ax_y, 2*qw + qgap, Inches(0.3),
     [[("WEARABLE / HANDS-FREE CAPTURE  →", 10, MUTE, True)]], align=PP_ALIGN.CENTER)
text(s, mx - Inches(0.42), my + qh - Inches(0.5), Inches(0.35), Inches(1.4),
     [[("A", 10, MUTE, True)], [("G", 10, MUTE, True)], [("E", 10, MUTE, True)],
      [("N", 10, MUTE, True)], [("T", 10, MUTE, True)], [("I", 10, MUTE, True)],
      [("C", 10, MUTE, True)]], align=PP_ALIGN.CENTER, space_after=0, line_spacing=0.9)

# right takeaways
bullet_list(s, [
    ("DAX / Abridge: no eyewear, no visual capture — and they stop at a note.",),
    ("Vuzix: best enterprise glasses — no agentic AI, aimed at factory work.",),
    ("AlphaAgent is the only wearable + agentic player that acts.", True),
    ("A scribe must become a hardware company to catch us; a hardware company must build our agent.",),
], Inches(8.0), Inches(2.25), Inches(4.6), size=12.5, gap=13)
footer(s, 6)
notes(s, "'No competition' is never true. Microsoft's DAX is a real, well-funded ambient scribe — but "
"it's a microphone in the room: no eyewear, never captures the surgical image, and stops at a draft "
"a human still acts on. Vuzix makes the best enterprise glasses but has no clinical agent and "
"targets factory and field work. We sit in the one quadrant neither can reach: wearable, visual, AND "
"agentic. To catch us a scribe must become a hardware company and a hardware company must build our "
"agent. That intersection is the moat.")

# ===========================================================================
# SLIDE 7 — Market
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Market Opportunity", "Land the OR workflow, expand to surgical intelligence")
tiers = [
    ("TAM", "$8T global healthcare — documentation & workflow automation", Inches(11.9), NAVY, SOFT),
    ("SAM", "US clinical documentation + surgical workflow", Inches(9.2), NAVY2, SOFT),
    ("SOM", "High-volume surgical specialties — ortho, general, plastics — US hospitals & ASCs", Inches(6.5), TEAL, WHITE),
]
yy = Inches(2.2)
for name, desc, w, fill, sub_c in tiers:
    card(s, MARGIN, yy, w, Inches(0.95), fill)
    text(s, MARGIN + Inches(0.35), yy + Inches(0.17), Inches(1.0), Inches(0.6),
         [[(name, 16, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, MARGIN + Inches(1.45), yy + Inches(0.12), w - Inches(1.8), Inches(0.75),
         [[(desc, 12.5, WHITE if fill is TEAL else sub_c)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    yy += Inches(1.15)
text(s, MARGIN, Inches(5.85), Inches(12.0), Inches(0.8),
     [[("Wedge → platform: ", 15, NAVY, True),
       ("own the moment of capture in the OR today → automated diagnostics on a proprietary multimodal surgical dataset tomorrow.", 15, INK)]],
     line_spacing=1.2)
footer(s, 7)
notes(s, "We start narrow and deep: high-volume surgical specialties where documentation pain is "
"worst and photos matter most. Once we own the moment of capture in the OR, we sit on a proprietary "
"multimodal surgical dataset that opens diagnostics and decision support in an $8T market. Land the "
"workflow, expand to intelligence.")

# ===========================================================================
# SLIDE 8 — Business model
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Business Model", "B2B SaaS — hardware is the wedge, software is the margin")
bullet_list(s, [
    ("Recurring per-surgeon subscription to hospitals, surgical groups, and ASCs.", True),
    ("Hardware is an enabler, not the margin — bundled at low markup to drive adoption.",),
    ("Land-and-expand: one service line → department → hospital system.", True),
    ("Sticky by design: ripping it out means going back to the keyboard.",),
], MARGIN, Inches(2.15), Inches(7.5), size=15, gap=15)
panel(s, "Land & expand", [
    "Hospital / ASC signs one service line",
    "↓   per-surgeon subscription",
    "Department-wide expansion",
    "↓   surgeon word of mouth",
    ("System-wide standard of care", True),
])
footer(s, 8)
notes(s, "We make money on software, not glasses. Hardware is the wedge into the OR; recurring "
"revenue is the per-surgeon subscription. We land in one service line, prove the time savings, and "
"expand. It's sticky — once a surgeon works hands-free, manual charting is unthinkable. That's "
"word-of-mouth growth inside a hospital.")

# ===========================================================================
# SLIDE 9 — Regulatory
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Regulatory & Path to Market", "Regulation as a moat — built in, not bolted on")
bullet_list(s, [
    ("Pathway: Class II 510(k) — predicate: surgical lighting / headwear.", True),
    ("Lean ISO 13485 QMS from day one; HIPAA + SOC 2 compliance-native.",),
    ("Q3 2026: FDA Pre-Submission locks the testing protocol.", True),
    ("Honest milestone: this round funds a cleared PATHWAY + live pilots.", True),
    ("Full 510(k) clearance completes ~2027, on the next round.", False, True),
], MARGIN, Inches(2.15), Inches(7.5), size=14.5, gap=13)
panel(s, "Timeline", [
    "NOW      HIPAA / SOC 2 posture",
    "  ↓        ISO 13485 lean QMS",
    "Q3 '26   FDA Pre-Submission",
    "  ↓        Hospital pilots",
    "2027     510(k) submission → clearance",
    ("Each step is a barrier to whoever chases us.", True),
])
footer(s, 9)
notes(s, "Regulation is a moat if you treat it as one. Class II 510(k) with a lean quality system "
"from day one; the Pre-Submission this quarter locks the testing protocol with the FDA. Straight "
"talk: this round gets us to a validated pathway and live pilots — full clearance lands in 2027 on "
"the next round. That's the exact inflection that de-risks the Series A, and it's why the regulatory "
"work is a barrier to everyone chasing us.")

# ===========================================================================
# SLIDE 10 — Team
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "Team", "A hardware-and-software product needs a hardware-and-software team")
tw = Inches(5.85)
for i, (nm, role, lines) in enumerate([
    ("Nathanael Gill", "CEO — Software & AI",
     ["Builds the agent: memory, tool-calling, FHIR integration.",
      "Healthcare analytics at HCA — clinical data pipelines & HIPAA from the inside.",
      "Sole salaried founder today; ships the pilot on off-the-shelf Vuzix."]),
    ("Manny Figueroa", "CTO — Hardware",
     ["Mechanical/systems engineer at Raytheon — reliability-critical hardware.",
      "Leads the custom enclosure & ruggedization phase.",
      "Joins equity-only; draws pay after the pilot proves the wedge."]),
]):
    x = MARGIN + i * (tw + Inches(0.25))
    card(s, x, Inches(2.15), tw, Inches(3.35), LIGHT, LINE)
    rect(s, x, Inches(2.15), tw, Inches(0.07), TEAL)
    text(s, x + Inches(0.35), Inches(2.5), tw - Inches(0.7), Inches(0.45),
         [[(nm, 17, NAVY, True)]])
    text(s, x + Inches(0.35), Inches(2.98), tw - Inches(0.7), Inches(0.35),
         [[(role, 12.5, TEAL, True)]])
    tb = s.shapes.add_textbox(x + Inches(0.35), Inches(3.45), tw - Inches(0.7), Inches(1.9))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for j, ln in enumerate(lines):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_after = Pt(8); p.line_spacing = 1.15
        r = p.add_run(); r.text = "▸  " + ln; _style(r, 12, INK)
text(s, MARGIN, Inches(5.8), Inches(12.0), Inches(0.6),
     [[("Capital-efficient by design: ", 14, NAVY, True),
       ("one salary, off-the-shelf hardware, custom build only after the pilot earns it.", 14, MUTE)]])
footer(s, 10)
notes(s, "This is a hardware-and-software product and we're a hardware-and-software team. I build the "
"agent and know the healthcare data pipelines and HIPAA rules from the inside at HCA. Manny builds "
"reliability-critical hardware at Raytheon — exactly the discipline you want when your device goes "
"into an OR. And we've sequenced it for capital efficiency: I'm the only salaried founder today, on "
"off-the-shelf Vuzix, and Manny comes on equity — drawing pay only once the pilot proves the wedge "
"and we start the custom hardware build. You need both to build this; few teams have both.")

# ===========================================================================
# SLIDE 11 — Ask & Use of Funds
# ===========================================================================
s = prs.slides.add_slide(BLANK)
header(s, "The Ask", "~$500k to reach multi-hospital pilots with the FDA machine underway")

tx, ty, tw_ = MARGIN, Inches(2.15), Inches(7.15)
rows = [
    ("USE OF FUNDS (midpoint burn)", "", True),
    ("Hardware & manufacturing — LX1 fleet + post-pilot tooling", "~$103k", False),
    ("Regulatory & compliance — QMS, HIPAA/SOC 2, 510(k) started", "~$204k", False),
    ("18-mo solo runway + Manny (post-pilot) + cloud", "~$219k", False),
    ("Total — Aggressive scenario", "~$525k", True),
]
rh = Inches(0.78)
for i, (label, val, strong) in enumerate(rows):
    y = ty + rh * i
    bg = NAVY if strong else (LIGHT if i % 2 else WHITE)
    rect(s, tx, y, tw_, rh, bg, LINE if not strong else None)
    text(s, tx + Inches(0.28), y, tw_ - Inches(1.9), rh,
         [[(label, 12.5, WHITE if strong else INK, strong)]],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    if val:
        text(s, tx + tw_ - Inches(1.55), y, Inches(1.3), rh,
             [[(val, 13.5, GOLD if strong else NAVY, True)]],
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

px, pw_ = Inches(8.3), Inches(4.33)
card(s, px, ty, pw_, rh * 5, NAVY)
rect(s, px, ty, pw_, Inches(0.06), GOLD)
ix, iw = px + Inches(0.35), pw_ - Inches(0.7)
text(s, ix, ty + Inches(0.3), iw, Inches(0.3), [[("THE ASK", 12, GOLD, True)]])
text(s, ix, ty + Inches(0.62), iw, Inches(0.85), [[("~$500k", 40, WHITE, True)]])
text(s, ix, ty + Inches(1.62), iw, Inches(0.75),
     [[("$300k = efficient floor — 12-mo runway, FDA Pre-Sub, single-site pilots.", 12, TEAL)]],
     line_spacing=1.15)
text(s, ix, ty + Inches(2.5), iw, Inches(1.2),
     [[("“We prove the wedge on $100k;", 14, WHITE, True, True)],
      [("we clear the regulatory moat on $500k.”", 14, WHITE, True, True)]],
     line_spacing=1.2, space_after=3)
footer(s, 11)
notes(s, "We're raising roughly $500k; here's exactly where it goes, bottoms-up. And note how "
"capital-efficient this is: I'm the only salaried founder — Manny joins on equity and only draws pay "
"AFTER the pilot proves the wedge, so the base runway carries one salary, not two. About $103k is "
"hardware — the LX1 fleet plus the post-pilot tooling. Roughly $204k is the regulatory moat — the "
"quality system, HIPAA and SOC 2, and getting the 510(k) underway. The rest is an 18-month solo "
"runway plus Manny's post-pilot onboarding and cloud. If you want the efficient version, $300k gets "
"a 12-month runway and single-site pilots on off-the-shelf Vuzix. And unlike most MedTech, we can "
"genuinely prove the wedge on ~$100k because the hardware spend stays deferred until it's earned.")

# ===========================================================================
# SLIDE 12 — Close
# ===========================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, 0, SW, Inches(0.09), TEAL)
text(s, MARGIN, Inches(1.9), Inches(11.9), Inches(1.0),
     [[("Own the surgeon's point of view.", 36, WHITE, True)]])
text(s, MARGIN, Inches(2.85), Inches(11.9), Inches(0.7),
     [[("Own the future of the operating room.", 24, TEAL, True)]])
rect(s, MARGIN, Inches(3.75), Inches(2.2), Pt(3), GOLD)
tb = s.shapes.add_textbox(MARGIN, Inches(4.15), Inches(11.5), Inches(1.8))
tf = tb.text_frame; tf.word_wrap = True
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
for j, ln in enumerate([
    "Today: kill documentation friction in the OR.",
    "Tomorrow: the intelligent capture-and-action layer for all of surgery.",
    "Always: building the proprietary surgical dataset that powers automated diagnostics.",
]):
    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
    p.space_after = Pt(12); p.line_spacing = 1.15
    r = p.add_run(); r.text = "▸  " + ln; _style(r, 15, WHITE)
text(s, MARGIN, Inches(6.55), Inches(11.9), Inches(0.4),
     [[("Nathanael Gill  ·  Manny Figueroa  ·  AlphaAgent Healthcare Solutions", 11, SOFT)]])
notes(s, "Documentation is the wedge, not the destination. Whoever owns the moment of capture in the "
"OR — the surgeon's point of view — owns the data layer the next decade of surgical AI is built on. "
"That's what we're building, and this round is how it starts.")

# ---------------------------------------------------------------------------
out = "AlphaAgent_Pitch_Deck.pptx"
prs.save(out)
print(f"Saved {out} with {len(prs.slides._sldIdLst)} slides.")
